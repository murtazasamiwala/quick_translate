"""Module to translate given set of documents"""
import os
from os.path import abspath
import win32com.client as win32
import xlrd
import pptx
import csv
import sys
from google.oauth2 import service_account
from google.cloud import translate_v2 as translate
base_path = os.path.dirname(abspath('__file__'))
allowed_ext = ['doc', 'docx', 'rtf', 'xls', 'xlsx', 'pptx', 'txt', 'csv']
if 'dont_delete_ignore' not in os.listdir(base_path):
    os.mkdir('dont_delete_ignore')
    kmsg_1 = 'Key not found.\n'
    kmsg_2 = 'Key folder has been created.\nSave key file to this folder.'
    fkmsg = kmsg_1 + kmsg_2
    result = open('script_result.txt', 'w', encoding='utf8')
    result.write(fkmsg)
    result.close()
    sys.exit()
if 'results_dir' not in os.listdir(base_path):
    os.mkdir('results_dir')
results_path = base_path + '\\' + 'results_dir'
key_folder = base_path + '\\' + 'dont_delete_ignore'
ignored_fol = ['result_dir', 'dont_delete_ignore']
key_path = key_folder + '\\' + os.listdir(key_folder)[0]
credentials = service_account.Credentials.from_service_account_file(
    key_path, scopes=["https://www.googleapis.com/auth/cloud-platform"])


def extract_text(fname, path=base_path):
    """Extract text from given document."""
    if fname.split('.')[-1] in ['doc', 'docx', 'rtf']:
        word = win32.gencache.EnsureDispatch('Word.Application')
        doc = word.Documents.Open(path+'\\'+fname)
        txt = doc.Content.Text
        doc.Close(False)
    elif fname.split('.')[-1] in ['xls', 'xlsx']:
        workbook = xlrd.open_workbook(path+'\\'+fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                    if new_output:
                        txt += ' '.join([str(i) for i in new_output]) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(path+'\\'+fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(path+'\\'+fname, 'r', encoding='utf8')
        txt = text_doc.read()
        text_doc.close()
    elif fname.endswith('.csv'):
        csv_doc = open(path+'\\'+fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    return txt, fname


def translate_text(doc, fname):
    """Check if text is already tranlsated. If not, translate it."""
    gt_out = None
    if gt_out is None:
        split_doc = [doc[:4000], doc[4000:]]
        translate_client = translate.Client(credentials=credentials)
        result = translate_client.translate(split_doc[0], target_language='en')
        split_doc[0] = result['translatedText']
        gt_out = ' '.join(split_doc)
    return gt_out


def save_files(gt_text, name):
    """Save translated text in separate folder as txt file."""
    word = win32.Dispatch('Word.Application')
    document = word.Documents.Add()
    document.Content.Text = gt_text
    document.SaveAs(results_path + '\\' + name.split('.')[0])
    document.Close(False)
    return


def report_file(outcome, fname):
    """Report whether script run successful or not."""
    if 'report_file.txt' not in os.listdir(base_path):
        report = open('report_file.txt', 'a', encoding='utf8')
        report.close()
    if outcome == 'success':
        msg = 'Translated successfuly: {}'.format(fname)
    elif outcome == 'error':
        msg_1 = 'ERROR in translating {}. '.format(fname)
        instruc = 'Cut paste the text into same doc and try again.'
        msg = msg_1 + instruc
    elif outcome == 'already':
        msg = '{} is already translated. Check results folder'.format(fname)
    report = open('report_file.txt', 'a', encoding='utf')
    report.write(msg + '\n')
    report.close()
    return



def folder_run(path=base_path):
    """Run script over all allowed files in folder."""
    for i in os.listdir(path):
        if i.endswith('zip'):
            pass
        elif i == 'report_file.txt':
            pass
        elif i in ignored_fol:
            pass
        elif i.split('.')[-1] in allowed_ext:
            src, trans_name = extract_text(i)
            if i.split('.')[0] + '.docx' in os.listdir(results_path):
                report_file('already', trans_name)
            else:
                g_translated = translate_text(src, trans_name)
                if g_translated is not None:
                    save_files(g_translated, trans_name)
                    report_file('success', trans_name)
                else:
                    report_file('error', trans_name)
        else:
            pass
    return


if __name__ == '__main__':
    folder_run()
