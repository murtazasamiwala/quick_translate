"""Module to translate given set of documents"""
import os
from os.path import abspath
from googletrans import Translator
import win32com.client as win32
import xlrd
import pptx
import csv
base_path = os.path.dirname(abspath('__file__'))
allowed_ext = ['doc', 'docx', 'rtf', 'xls', 'xlsx', 'pptx', 'txt', 'csv']


def extract_text(fname, path=base_path):
    """Extract text from given document."""
    if fname.split('.')[-1] in ['doc', 'docx', 'rtf']:
        word = win32.Dispatch('Word.Application')
        doc = word.Documents.Open(path+'\\'+fname)
        txt = doc.Content.Text
        doc.Close(False)
        word.Quit()
    elif fname.split('.')[-1] in ['xls', 'xlsx']:
        workbook = xlrd.open_workbook(fname)
        sheets_name = workbook.sheet_names()
        txt = '\n'
        for names in sheets_name:
            worksheet = workbook.sheet_by_name(names)
            num_rows = worksheet.nrows
            num_cells = worksheet.ncols
            for curr_row in range(num_rows):
                new_output = []
                for index_col in range(num_cells):
                    value = worksheet.cell_value(curr_row, index_col)
                    if value:
                        new_output.append(value)
                    if new_output:
                        txt += ' '.join(new_output) + '\n'
    elif fname.endswith('.pptx'):
        presentation = pptx.Presentation(fname)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        txt = '\n\n'.join(text_runs)
    elif fname.endswith('.txt'):
        text_doc = open(fname, 'r', encoding='utf8')
        txt = text_doc.read()
    elif fname.endswith('.csv'):
        csv_doc = open(fname, 'r', encoding='utf8')
        csv_reader = csv.reader(csv_doc, delimiter=',')
        txt = '\n'.join(['\t'.join(row) for row in csv_reader])
    return txt, fname


def doc_split(doc):
    """Split text into small chunks readable by google translate."""
    translator = Translator()
    lan = translator.detect(doc[:1000]).lang
    if lan in ['ko', 'pt']:
        tokens = doc.split('.')
        tokens = [i + '.' for i in tokens]
    elif lan in ['ja', 'zh-CN', 'zh-TW']:
        tokens = doc.split('。')
        tokens = [i + '。' for i in tokens]
    else:
        tokens = doc.split('.')
        tokens = [i + '.' for i in tokens]
    split = []
    len_counter = 0
    temp_list = []
    final = []
    for i in range(len(tokens)):
        if len_counter + len(tokens[i]) + len(temp_list) - 1 < 1800:
            len_counter = len_counter + len(tokens[i])
            temp_list.append(tokens[i])
        else:
            len_counter = len(tokens[i])
            split.append(temp_list)
            temp_list = []
            temp_list.append(tokens[i])
    split.append(temp_list)
    final = [''.join(i) for i in split]
    return final


def translate_text(src_list, fname):
    """Translate text given as list of strings."""
    google_output = []
    for i in src_list:
        try:
            translator = Translator()
            google_output.append(translator.translate(i, dest='en').text)
        except ValueError:
            return None
    translated = ' '.join(google_output)
    return translated


def save_files(gt_text, name):
    """Save translated text in separate folder as txt file."""
    if 'translations_dir' not in os.listdir(base_path):
        os.mkdir('translations_dir')
    save_dir = base_path + '\\' + 'translations_dir' + '\\' + 'translated_'
    txt_file = open(save_dir + '{}.txt'.format(name), 'w', encoding='utf8')
    txt_file.write(gt_text)
    txt_file.close()
    return


def report_file(outcome, fname):
    """Report whether script run successful or not."""
    if 'report_file.txt' not in os.listdir(base_path):
        report = open('report_file.txt', 'a', encoding='utf8')
        report.close()
    if outcome == 'success':
        msg = 'Translated successfuly: {}'.format(fname)
    elif outcome == 'error':
        msg_1 = 'ERROR in translating {}. '.format(fname)
        instruc = 'Cut paste the text into same doc and try again.'
        msg = msg_1 + instruc
    report = open('report_file.txt', 'a', encoding='utf')
    report.write(msg + '\n')
    report.close()
    return



def folder_run(path=base_path):
    """Run script over all allowed files in folder."""
    for i in os.listdir(path):
        if (os.path.isdir(i)) | (i.endswith('zip')):
            pass
        elif i == 'report_file.txt':
            pass
        elif i.split('.')[-1] in allowed_ext:
            src, file_name = extract_text(i)
            source_list = doc_split(src)
            g_translated = translate_text(source_list, file_name)
            if g_translated is not None:
                save_files(g_translated, file_name)
                report_file('success', file_name)
            else:
                report_file('error', file_name)
        else:
            pass
    return


if __name__ == '__main__':
    folder_run()
